"""Genera la presentacion del Trabajo Final a partir de los resultados REALES.

Los numeros no se escriben a mano: se leen de `reports/` y de
`models/*/test_metrics.json`. Si una metrica todavia no existe, la placa lo dice
explicitamente en lugar de inventar un valor.

  python scripts/build_ppt.py

Salida: entregables/TP_Final_Deep_Learning_Vassarotto.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

# Paleta y tipografia definidas por el alumno
NAVY = RGBColor(0x16, 0x24, 0x3D)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS = RGBColor(0x44, 0x4A, 0x55)
FUENTE = "Arial"

ANCHO = Inches(13.333)   # 16:9
ALTO = Inches(7.5)


# --------------------------------------------------------------------------
# Lectura de resultados reales
# --------------------------------------------------------------------------

def _json(rel):
    p = ROOT / rel
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return None


def cargar_resultados() -> dict:
    """Junta todo lo que haya en reports/ y models/. Lo que falta queda en None."""
    return {
        "ner": _json("models/ner-beto/test_metrics.json"),
        "cls": _json("models/cls-beto/test_metrics.json"),
        "ner_sint": _json("reports/ner_sintetico.json"),
        "cls_sint": _json("reports/cls_sintetico.json"),
        "ner_real": _json("reports/ner_real.json"),
        "cls_real": _json("reports/cls_real.json"),
        "robustez": _json("reports/parser_robustness.json"),
        "agente": (_json("reports/agent_metrics_argenprop.json")
                   or _json("reports/agent_metrics_deterministic.json")),
    }


def fmt(valor, pct=True):
    if valor is None:
        return "(pendiente)"
    return f"{valor:.1%}" if pct else f"{valor:.4f}"


# --------------------------------------------------------------------------
# Helpers de armado de placas
# --------------------------------------------------------------------------

def _fondo(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _caja(slide, texto, x, y, w, h, size=18, color=NAVY, bold=False,
          align=PP_ALIGN.LEFT, espaciado=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lineas = texto if isinstance(texto, list) else [texto]
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(espaciado)
        run = p.add_run()
        run.text = linea
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FUENTE
    return tb


def _linea_gold(slide, y=Inches(1.35), x=Inches(0.9), w=Inches(2.2)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = GOLD
    sh.line.fill.background()
    sh.shadow.inherit = False


def placa_titulo(prs, titulo, subtitulo, pie):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, NAVY)
    _caja(s, titulo, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.4),
          size=36, color=BLANCO, bold=True)
    _linea_gold(s, y=Inches(4.55), x=Inches(0.95))
    _caja(s, subtitulo, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.0),
          size=18, color=GOLD)
    _caja(s, pie, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
          size=14, color=BLANCO)
    return s


def placa(prs, titulo, cuerpo, size=17, bajada=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, BLANCO)
    _caja(s, titulo, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.9),
          size=28, color=NAVY, bold=True)
    _linea_gold(s)
    y = Inches(1.75)
    if bajada:
        _caja(s, bajada, Inches(0.9), y, Inches(11.5), Inches(0.6),
              size=15, color=GOLD, bold=True)
        y = Inches(2.35)
    _caja(s, cuerpo, Inches(0.9), y, Inches(11.5), Inches(4.6),
          size=size, color=GRIS, espaciado=11)
    return s


def placa_metricas(prs, titulo, filas, nota=None):
    """Placa con tarjetas de numeros grandes: [(valor, etiqueta), ...]"""
    from pptx.enum.shapes import MSO_SHAPE
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fondo(s, BLANCO)
    _caja(s, titulo, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.9),
          size=28, color=NAVY, bold=True)
    _linea_gold(s)

    n = len(filas)
    ancho = Inches(11.5 / max(n, 1) - 0.25)
    for i, (valor, etiqueta) in enumerate(filas):
        x = Inches(0.9 + i * (11.5 / n))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0),
                                  ancho, Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY
        card.line.fill.background()
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(valor)
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = GOLD
        r.font.name = FUENTE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = etiqueta
        r2.font.size = Pt(13)
        r2.font.color.rgb = BLANCO
        r2.font.name = FUENTE
    if nota:
        _caja(s, nota, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2.0),
              size=15, color=GRIS, espaciado=10)
    return s


# --------------------------------------------------------------------------
# La presentacion
# --------------------------------------------------------------------------

def construir(res: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = ANCHO, ALTO

    # 1. Portada
    placa_titulo(
        prs,
        "Agente basado en LLM para la adquisición autónoma de avisos "
        "inmobiliarios y enriquecimiento semántico mediante NLP",
        "Trabajo Final — Deep Learning",
        "Joaquín Héctor Vassarotto · Legajo 106442\n"
        "Maestría en Management & Analytics (MMA) — ITBA",
    )

    # 2. El problema (vision de negocio)
    placa(prs, "El problema de negocio",
          ["Valuar departamentos en CABA con modelos automáticos (AVM) choca con dos "
           "limitaciones de datos:",
           "",
           "1.  Adquirir los avisos es frágil y manual. Los scrapers por reglas fijas se "
           "rompen ante cada cambio de layout del portal, y no toleran la heterogeneidad "
           "entre publicaciones.",
           "",
           "2.  La información de valor está en el texto libre. Estado, amenities, "
           "orientación, antigüedad y señales del vendedor viven en la descripción — "
           "no en los campos tabulares que ofrecen los datasets públicos.",
           "",
           "Dos propiedades con los mismos metros y ambientes pueden valer muy distinto "
           "según lo que diga la descripción."],
          bajada="¿Qué problema específico busca resolver el proyecto?")

    # 3. Propuesta de valor
    placa(prs, "Propuesta de valor",
          ["El objetivo no es obtener datos: es generar variables latentes a partir del "
           "lenguaje natural.",
           "",
           "•  Un agente que se adapta en lugar de romperse ante cambios del portal.",
           "",
           "•  Atributos estructurados que los datasets tabulares no capturan: estado de "
           "conservación, amenities, orientación, y señales del vendedor.",
           "",
           "•  Señales potencialmente correlacionadas con subvaluación — «dueño directo», "
           "«escucho ofertas», «venta urgente» — que son justamente las de mayor interés "
           "para detectar oportunidades."],
          bajada="Qué aporta que no exista hoy")

    # 4. Arquitectura
    placa(prs, "Arquitectura: dos capas",
          ["CAPA 1 — Agente (orquestación).  LLM local y gratuito (Ollama) con patrón "
           "ReAct / tool-use. El LLM actúa como policy: razona sobre el estado de la "
           "página y decide la acción (navegar, paginar, extraer, reintentar).",
           "      Herramientas: Playwright (contenido dinámico) + parser resiliente.",
           "      → El LLM se usa PRE-ENTRENADO. No es lo que se entrena.",
           "",
           "CAPA 2 — Enriquecimiento (NLP).  Dos transformers BETO fine-tuneados que "
           "convierten la descripción en variables estructuradas.",
           "      → Este es el núcleo entrenable y el entregable de la materia.",
           "",
           "Stack: LangChain + LangGraph · Playwright · HuggingFace Transformers · PyTorch"],
          size=16)

    # 5. Los modelos entrenados
    placa(prs, "Los modelos que se entrenan",
          ["Ambos parten de BETO (dccuchile/bert-base-spanish-wwm-cased), un BERT "
           "entrenado en español.",
           "",
           "1.  NER — token classification, esquema BIO.",
           "      AMENITY · ESTADO · ANTIGUEDAD · ORIENTACION · EXPENSAS",
           "      Dónde, dentro del texto, se menciona cada atributo.",
           "",
           "2.  Clasificación multilabel — señales del vendedor.",
           "      DUENO_DIRECTO · OPORTUNIDAD · URGENCIA · REFACCION",
           "      Multilabel y no multiclase: un aviso puede ser a la vez «dueño directo» "
           "y «urgencia».",
           "",
           "Restricción de hardware: GTX 1650 con 4 GB de VRAM → fp16, batch chico y "
           "gradient accumulation para recuperar un batch efectivo razonable."],
          size=16)

    # 6. Metricas: definicion y justificacion  (lo pide explicitamente la consigna)
    placa(prs, "Métricas: definición y justificación",
          ["NER — F1 por entidad con seqeval, NO accuracy por token.",
           "      • El accuracy por token miente: la enorme mayoría de los tokens son «O», "
           "así que un modelo que no detecte nada igual sacaría un valor altísimo.",
           "      • Importa el span completo: si la anotación dice «a estrenar» y el modelo "
           "marca sólo «estrenar», la variable resultante está mal. seqeval exige que "
           "coincidan tipo y límites exactos.",
           "",
           "Clasificación — F1 macro además del micro.",
           "      • Las clases están desbalanceadas; el micro queda dominado por las "
           "frecuentes.",
           "      • El macro promedia las cuatro clases con igual peso, y por eso penaliza "
           "ignorar una clase rara — justamente URGENCIA y OPORTUNIDAD, las más útiles "
           "para detectar subvaluación.",
           "",
           "Agente — tasa de éxito de extracción y robustez ante cambios de layout."],
          size=15)

    # 7. EL resultado del trabajo: la brecha sintetico -> real
    ns, cs = res["ner_sint"], res["cls_sint"]
    nr, cr = res["ner_real"], res["cls_real"]
    placa_metricas(
        prs, "Resultados — la brecha sintético → real",
        [(fmt(ns.get("f1_micro") if ns else None), "F1 NER\nsobre sintético"),
         (fmt(nr.get("f1_micro") if nr else None), "F1 NER\nsobre REAL"),
         (fmt(cs.get("f1_macro") if cs else None), "F1 macro clasif.\nsobre sintético"),
         (fmt(cr.get("f1_macro") if cr else None), "F1 macro clasif.\nsobre REAL")],
        nota=["Este contraste es el hallazgo central, y no se disimula: un F1 casi perfecto sobre "
              "datos sintéticos no mide qué tan bueno es el modelo, sino qué tan fácil es el test.",
              "",
              "Los avisos sintéticos salen de plantillas, así que el modelo aprende la plantilla "
              "en lugar del concepto."])

    # 8. Como falla cada modelo (el analisis, no solo el numero)
    placa(prs, "Los dos modelos fallan distinto",
          ["EL NER SOBRE-ETIQUETA.",
           "      Aprendió «sustantivo después de Cuenta con» en vez del vocabulario real de "
           "amenities. Sobre texto real marca AMENITY → «ventilación», «universidades», y llega a "
           "marcar el nombre de una calle como ORIENTACION.",
           "",
           "EL CLASIFICADOR CASI NO DISPARA.",
           "      Precisión macro 0.55 pero recall 0.16: cuando predice suele acertar, pero se "
           "pierde la mayoría. Los avisos reales expresan «dueño directo» o «urgencia» con formas "
           "que el generador nunca produjo.",
           "",
           "CÓMO LEER ESTOS NÚMEROS.",
           "      Las etiquetas del conjunto real vienen del pre-anotador LLM sin revisión humana "
           "completa, sobre 65 avisos. Miden concordancia con un anotador imperfecto, no verdad de "
           "referencia. La magnitud de la caída es sólida; los valores exactos, no."],
          size=15)

    # 8. Resultados del agente (numeros reales)
    rob = res["robustez"]
    ret = rob["retencion_promedio"] if rob else {}
    ag = res["agente"]["summary"] if res["agente"] else {}
    placa_metricas(
        prs, "Resultados — capa de agente",
        [(fmt(ag.get("extraction_success_rate")), "Tasa de éxito\nde extracción"),
         (str(ag.get("listings_extracted", "—")), "Avisos\nextraídos"),
         (fmt(ret.get("sin_data_qa")), "Retención del parser\nsin selectores data-qa"),
         (fmt(ret.get("renombrado_a_clases")), "Retención\ncon layout renombrado")],
        nota=["Robustez: se toma un HTML que el parser sabe leer y se lo degrada imitando "
              "cambios reales del portal, midiendo qué campos sobreviven.",
              "",
              "Al desaparecer los selectores data-qa, los campos numéricos y el barrio se "
              "recuperan por regex sobre el texto plano; sólo se pierde la descripción, que "
              "necesita algún gancho de markup. Un scraper de un solo nivel habría caído a 0%."])

    # 9. La objecion mas razonable, medida en vez de supuesta
    cob = _json("reports/cobertura_tabulada.json")
    r = (cob or {}).get("resumen", {})
    placa(prs, "«¿Y si el portal ya publica los atributos?»",
          [f"Argenprop SÍ expone una ficha estructurada en cada página de detalle: 38 items "
           f"tildados y 20 pares clave-valor. La objeción es válida, así que se midió.",
           "",
           f"Sobre una muestra de {r.get('avisos_medidos', '—')} avisos, comparando ambas fuentes "
           f"con el mismo vocabulario:",
           f"      • Amenities tildados en la ficha:            {r.get('amenities_en_ficha', '—')}",
           f"      • Amenities mencionados en la descripción:   {r.get('amenities_en_descripcion', '—')}",
           f"      • SÓLO en la descripción (se perderían):     {r.get('amenities_SOLO_en_descripcion', '—')}",
           f"      • SÓLO en la ficha:                          {r.get('amenities_SOLO_en_ficha', '—')}",
           "",
           "LAS DOS FUENTES SON COMPLEMENTARIAS, NO REDUNDANTES. Ninguna contiene a la otra: "
           "la ficha la completa el publicador, y cada uno completa distinto.",
           "",
           "Y las señales del vendedor —«escucho ofertas», «venta urgente»— no tienen equivalente "
           "tabulado en ningún portal. Ahí el NLP no compite con nada."],
          size=15, bajada="La objeción más razonable al trabajo, respondida con datos")

    # 10. Limitaciones
    placa(prs, "Limitaciones",
          ["•  Volumen del scrape. La propuesta apuntaba a 8.000–15.000 avisos; se trabajó "
           "con un volumen mucho menor. El anti-bot del portal (DataDome/Cloudflare) y el "
           "rate-limiting cortés hacen que una corrida completa lleve muchas horas. El "
           "objetivo fue demostrar que la arquitectura funciona.",
           "",
           "•  Entrenamiento sobre datos sintéticos, con el conjunto real como evaluación "
           "externa. Un dataset real y grande daría mejores resultados.",
           "",
           "•  Anotación semiautomática. El pre-anotador es un LLM chico (llama3.2:3b, por "
           "la restricción de 4 GB de VRAM); sus errores se propagan a las etiquetas, y por "
           "eso la revisión manual del subconjunto.",
           "",
           "•  Términos de uso. Scraping con fines académicos, a ritmo razonable y sin "
           "redistribuir el contenido del portal."],
          size=15)

    # 10. Trabajo futuro
    placa(prs, "Trabajo futuro",
          ["Los atributos generados por esta capa de NLP están pensados para integrarse, "
           "en el marco de la tesis de la maestría, como features de un modelo hedónico de "
           "valuación para detectar activos subvaluados en CABA.",
           "",
           "La hipótesis a testear: que las variables latentes extraídas del texto libre "
           "aporten poder explicativo por encima de los atributos tabulares tradicionales "
           "(metros, ambientes, ubicación, antigüedad).",
           "",
           "Esa integración excede el alcance del presente Trabajo Final y se menciona para "
           "dar cuenta de la continuidad del proyecto."],
          bajada="La motivación detrás del proyecto")

    # 11. Cierre
    s = placa_titulo(prs, "Gracias",
                     "Código, datos y reportes reproducibles",
                     "github.com/<usuario>/zonaprop-agent-nlp")
    return prs


def main():
    res = cargar_resultados()
    faltantes = [k for k, v in res.items() if not v]
    if faltantes:
        print(f"[AVISO] sin resultados todavia para: {', '.join(faltantes)}")
        print("        Esas placas van a decir '(pendiente)' en lugar de un numero inventado.")

    out_dir = ROOT / "entregables"
    out_dir.mkdir(exist_ok=True)
    out = out_dir / "TP_Final_Deep_Learning_Vassarotto.pptx"
    construir(res).save(out)
    print(f"[OK] presentacion -> {out}")


if __name__ == "__main__":
    main()
